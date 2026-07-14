"""
app/rag.py — RAG Pipeline
==========================
Handles document ingestion, chunking, embedding via Ollama (nomic-embed-text),
and semantic retrieval using ChromaDB as the vector store.
"""

import asyncio, logging, re
from pathlib import Path
from typing import Any

import chromadb
from chromadb.config import Settings

logger = logging.getLogger(__name__)

# ── Optional document loaders ─────────────────────────────────────────────────
try:
    import pdfplumber
    _PDF_OK = True
except ImportError:
    _PDF_OK = False

try:
    from docx import Document as DocxDocument
    _DOCX_OK = True
except ImportError:
    _DOCX_OK = False

try:
    import openpyxl
    _XLSX_OK = True
except ImportError:
    _XLSX_OK = False

try:
    from pptx import Presentation
    _PPTX_OK = True
except ImportError:
    _PPTX_OK = False

try:
    import ollama as _ollama_sdk
    _OLLAMA_SDK = True
except ImportError:
    _OLLAMA_SDK = False


class RAGPipeline:
    EMBED_MODEL   = "nomic-embed-text"
    CHUNK_SIZE    = 500   # characters
    CHUNK_OVERLAP = 100   # characters
    COLLECTION    = "ragdocs"
    DB_PATH       = "./chroma_db"

    def __init__(self):
        self.client = chromadb.PersistentClient(
            path=self.DB_PATH,
            settings=Settings(anonymized_telemetry=False),
        )
        self.col = self.client.get_or_create_collection(
            name=self.COLLECTION,
            metadata={"hnsw:space": "cosine"},
        )
        logger.info("ChromaDB ready: %d documents indexed", self.col.count())

    # ── Ingestion ──────────────────────────────────────────────────────────────

    async def ingest(self, file_path: str, doc_id: str, filename: str) -> dict:
        """Load → chunk → embed → store. Returns stats dict."""
        text = self._load_file(file_path)
        chunks = self._chunk(text)

        logger.info("Ingesting '%s': %d chunks", filename, len(chunks))

        ids        = [f"{doc_id}_{i}" for i in range(len(chunks))]
        embeddings = await self._embed_batch(chunks)
        metadatas  = [{"doc_id": doc_id, "filename": filename, "chunk_idx": i}
                      for i in range(len(chunks))]

        self.col.add(
            ids=ids,
            embeddings=embeddings,
            documents=chunks,
            metadatas=metadatas,
        )
        return {"chunks": len(chunks), "chars": len(text)}

    # ── Retrieval ──────────────────────────────────────────────────────────────

    async def query(self, question: str, k: int = 3) -> dict:
        """Embed the question, retrieve top-k chunks."""
        if self.col.count() == 0:
            return {"chunks": [], "sources": []}

        q_embed = (await self._embed_batch([question]))[0]
        results = self.col.query(
            query_embeddings=[q_embed],
            n_results=min(k, self.col.count()),
            include=["documents", "metadatas", "distances"],
        )

        docs      = results["documents"][0]
        metas     = results["metadatas"][0]
        distances = results["distances"][0]

        chunks = [
            {
                "text":     doc,
                "source":   m["filename"],
                "score":    round(1 - dist, 4),
                "chunk_idx": m["chunk_idx"],
            }
            for doc, m, dist in zip(docs, metas, distances)
        ]
        sources = sorted(set(c["source"] for c in chunks))
        return {"chunks": chunks, "sources": sources}

    # ── Knowledge-base admin ───────────────────────────────────────────────────

    async def list_documents(self) -> list[dict]:
        """Return one record per original document (unique doc_ids)."""
        if self.col.count() == 0:
            return []
        results = self.col.get(include=["metadatas"])
        seen: dict[str, dict] = {}
        for m in results["metadatas"]:
            did = m["doc_id"]
            if did not in seen:
                seen[did] = {"doc_id": did, "filename": m["filename"], "chunks": 0}
            seen[did]["chunks"] += 1
        return list(seen.values())

    async def delete_document(self, doc_id: str):
        """Remove all chunks for a document from ChromaDB."""
        results = self.col.get(where={"doc_id": doc_id})
        if results["ids"]:
            self.col.delete(ids=results["ids"])
            logger.info("Deleted %d chunks for doc_id=%s", len(results["ids"]), doc_id)

    # ── Private helpers ────────────────────────────────────────────────────────

    def _load_file(self, path: str) -> str:
        ext = Path(path).suffix.lower()
        if ext == ".pdf":
            return self._load_pdf(path)
        if ext == ".docx":
            return self._load_docx(path)
        if ext in (".xlsx", ".xls"):
            return self._load_xlsx(path)
        if ext in (".pptx", ".ppt"):
            return self._load_pptx(path)
        # Plain text formats (.txt, .md, .json, .csv)
        return Path(path).read_text(encoding="utf-8", errors="replace")

    def _load_pdf(self, path: str) -> str:
        if not _PDF_OK:
            raise ImportError("Install pdfplumber: pip install pdfplumber")
        with pdfplumber.open(path) as pdf:
            return "\n\n".join(
                page.extract_text() or "" for page in pdf.pages
            )

    def _load_docx(self, path: str) -> str:
        if not _DOCX_OK:
            raise ImportError("Install python-docx: pip install python-docx")
        doc = DocxDocument(path)
        return "\n".join(p.text for p in doc.paragraphs)

    def _load_xlsx(self, path: str) -> str:
        if not _XLSX_OK:
            raise ImportError("Install openpyxl: pip install openpyxl")
        wb = openpyxl.load_workbook(path, data_only=True)
        out = []
        for sheet in wb.worksheets:
            out.append(f"[Sheet: {sheet.title}]")
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    out.append(" | ".join(cells))
        return "\n".join(out)

    def _load_pptx(self, path: str) -> str:
        if not _PPTX_OK:
            raise ImportError("Install python-pptx: pip install python-pptx")
        prs = Presentation(path)
        out = []
        for i, slide in enumerate(prs.slides, start=1):
            out.append(f"[Slide {i}]")
            for shape in slide.shapes:
                # Text in shapes/text boxes/titles
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = "".join(run.text for run in para.runs)
                        if text.strip():
                            out.append(text.strip())
                # Text inside tables
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text for cell in row.cells if cell.text]
                        if cells:
                            out.append(" | ".join(cells))
                # Speaker notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
                out.append("[Notes] " + slide.notes_slide.notes_text_frame.text.strip())
        return "\n".join(out)

    def _chunk(self, text: str) -> list[str]:
        """Sliding-window character chunker with sentence-boundary snapping."""
        text = re.sub(r"\s+", " ", text).strip()
        chunks, start = [], 0
        while start < len(text):
            end = min(start + self.CHUNK_SIZE, len(text))
            # Snap to sentence boundary
            snap = text.rfind(". ", start, end)
            if snap != -1 and snap > start + self.CHUNK_OVERLAP:
                end = snap + 1
            chunk = text[start:end].strip()
            if chunk:
                chunks.append(chunk)
            start += self.CHUNK_SIZE - self.CHUNK_OVERLAP
        return chunks

    async def _embed_batch(self, texts: list[str]) -> list[list[float]]:
        """Generate embeddings via Ollama SDK or HTTP fallback."""
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(None, self._embed_sync, texts)

    def _embed_sync(self, texts: list[str]) -> list[list[float]]:
        if _OLLAMA_SDK:
            return [
                _ollama_sdk.embeddings(model=self.EMBED_MODEL, prompt=t)["embedding"]
                for t in texts
            ]
        # HTTP fallback
        import httpx
        embeddings = []
        for t in texts:
            r = httpx.post(
                "http://localhost:11434/api/embeddings",
                json={"model": self.EMBED_MODEL, "prompt": t},
                timeout=30,
            )
            r.raise_for_status()
            embeddings.append(r.json()["embedding"])
        return embeddings
